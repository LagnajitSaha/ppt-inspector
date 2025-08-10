import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .models import SlideContent

class ContentExtractor:
    def __init__(self, temp_dir="tmp"):
        self.temp_dir = temp_dir
        os.makedirs(temp_dir, exist_ok=True)

    def extract(self, pptx_path: str):
        prs = Presentation(pptx_path)
        slides_data = []

        for idx, slide in enumerate(prs.slides, start=1):
            text_content = []
            tables_content = []
            images_content = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_content.append(shape.text.strip())

                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_data = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    tables_content.append(table_data)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_path = os.path.join(self.temp_dir, f"slide_{idx}_{len(images_content)}.png")
                    with open(image_path, "wb") as f:
                        f.write(shape.image.blob)
                    images_content.append(image_path)

            slides_data.append(
                SlideContent(
                    slide_number=idx,
                    text=text_content,
                    tables=tables_content,
                    images=images_content
                )
            )

        return slides_data
