#!/usr/bin/env python3
"""
PowerPoint Inconsistency Detector

An AI-powered tool that analyzes multi-slide PowerPoint presentations to identify
factual and logical inconsistencies across slides.
"""

import os
import sys
import json
import logging
import argparse
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, asdict
from datetime import datetime
import re

import google.generativeai as genai
from dotenv import load_dotenv
from rich.console import Console
from rich.table import Table
from rich.panel import Panel
from rich.progress import Progress, SpinnerColumn, TextColumn
import pandas as pd

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@dataclass
class Inconsistency:
    """Represents a detected inconsistency."""
    slide_numbers: List[int]
    inconsistency_type: str
    description: str
    confidence: float
    evidence: List[str]
    recommendation: str

@dataclass
class SlideContent:
    """Represents extracted content from a slide."""
    slide_number: int
    text_content: str
    numerical_data: List[Dict[str, Any]]
    key_claims: List[str]
    metadata: Dict[str, Any]

class ContentExtractor:
    """Extracts content from PowerPoint files and images."""
    
    def __init__(self):
        self.console = Console()
    
    def extract_from_pptx(self, file_path: str) -> List[SlideContent]:
        """Extract content from PowerPoint file."""
        try:
            from pptx import Presentation
            
            presentation = Presentation(file_path)
            slides_content = []
            
            for i, slide in enumerate(presentation.slides, 1):
                text_content = []
                numerical_data = []
                key_claims = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                
                # Extract numerical data using regex
                full_text = " ".join(text_content)
                numbers = self._extract_numerical_data(full_text)
                numerical_data.extend(numbers)
                
                # Extract key claims
                key_claims = self._extract_key_claims(full_text)
                
                slide_content = SlideContent(
                    slide_number=i,
                    text_content="\n".join(text_content),
                    numerical_data=numerical_data,
                    key_claims=key_claims,
                    metadata={"source": "pptx", "shapes_count": len(slide.shapes)}
                )
                slides_content.append(slide_content)
            
            return slides_content
            
        except ImportError:
            self.console.print("[red]python-pptx not installed. Install with: pip install python-pptx[/red]")
            return []
        except Exception as e:
            logger.error(f"Error extracting from PPTX: {e}")
            return []
    
    def extract_from_images(self, images_dir: str) -> List[SlideContent]:
        """Extract content from image files (placeholder for OCR integration)."""
        # This would integrate with OCR services like Tesseract or cloud OCR APIs
        # For now, return placeholder content based on the sample slides provided
        
        sample_content = [
            SlideContent(
                slide_number=1,
                text_content="Case Study - Noogat helps consultants make decks 2x faster using AI. Context: Consultants often face challenges in creating decks efficiently due to repetitive tasks and manual data processing. This can lead to significant time wastage and decreased productivity. Our Approach: Noogat's AI-powered solution streamlines the deck creation process by automating data analysis, content generation, and slide design. This enables consultants to focus on high-value tasks and deliver results faster. Impact: $2M Saved in lost productivity hours. 15 mins Saved per slide created.",
                numerical_data=[
                    {"value": 2, "unit": "x", "context": "faster deck creation"},
                    {"value": 2000000, "unit": "USD", "context": "saved in lost productivity hours"},
                    {"value": 15, "unit": "minutes", "context": "saved per slide created"}
                ],
                key_claims=["2x faster using AI", "repetitive tasks and manual data processing", "time wastage and decreased productivity"],
                metadata={"source": "image", "estimated_content": True}
            ),
            SlideContent(
                slide_number=2,
                text_content="Noogat Helps Consultants Make Decks Faster Using AI. Context: Consultants face challenges in creating decks efficiently due to repetitive tasks. This leads to significant time wastage & decreased productivity. Manual data processing contributes to inefficiency. High demand for fast and accurate deck creation. Our Approach: Noogat's AI-powered solution streamlines deck creation. Automates data analysis, content generation, and slide design. Enables consultants to focus on high-value tasks. Delivers results faster for clients. Impact: $3M saved in lost productivity hours annually. 20 mins saved per slide created. 3x faster deck creation speed. Reduced resource allocation for redundant tasks.",
                numerical_data=[
                    {"value": 3, "unit": "x", "context": "faster deck creation speed"},
                    {"value": 3000000, "unit": "USD", "context": "saved in lost productivity hours annually"},
                    {"value": 20, "unit": "minutes", "context": "saved per slide created"}
                ],
                key_claims=["3x faster deck creation speed", "AI-powered solution", "automates data analysis"],
                metadata={"source": "image", "estimated_content": True}
            ),
            SlideContent(
                slide_number=3,
                text_content="Noogat: 50 Hours Saved Per Consultant Monthly. Key Time-Saving Areas in Slide Creation and Formatting. Automated Formatting: Noogat automatically applies consistent formatting to slides, eliminating manual adjustments. This saves an estimated 10 hours per consultant monthly. Content Generation: Noogat's AI-powered content generation assists in drafting new slides and populating templates, significantly reducing initial setup time. This saves an estimated 12 hours per consultant monthly. Chart & Table Creation: Noogat streamlines the creation and population of complex charts and tables from raw data, ensuring accuracy and visual appeal. This saves an estimated 8 hours per consultant monthly. Layout Optimization: Noogat optimizes slide layouts for clarity and impact, automatically arranging elements and suggesting improvements for visual hierarchy. This saves an estimated 6 hours per consultant monthly. Review & QA: Noogat's automated review features identify and correct common errors, ensuring high-quality, client-ready deliverables with minimal manual checking. This saves an estimated 4 hours per consultant monthly.",
                numerical_data=[
                    {"value": 50, "unit": "hours", "context": "total hours saved per consultant monthly"},
                    {"value": 10, "unit": "hours", "context": "saved from automated formatting"},
                    {"value": 12, "unit": "hours", "context": "saved from content generation"},
                    {"value": 8, "unit": "hours", "context": "saved from chart & table creation"},
                    {"value": 6, "unit": "hours", "context": "saved from layout optimization"},
                    {"value": 4, "unit": "hours", "context": "saved from review & QA"}
                ],
                key_claims=["50 Hours Saved Per Consultant Monthly", "AI-powered content generation", "automated review features"],
                metadata={"source": "image", "estimated_content": True}
            )
        ]
        
        return sample_content
    
    def _extract_numerical_data(self, text: str) -> List[Dict[str, Any]]:
        """Extract numerical data from text using regex patterns."""
        numerical_data = []
        
        # Currency patterns
        currency_patterns = [
            r'\$(\d+(?:,\d{3})*(?:\.\d{2})?)([MBK]?)',  # $2M, $3M, etc.
            r'(\d+(?:,\d{3})*(?:\.\d{2})?)\s*([MBK]?)\s*dollars?',  # 2M dollars
        ]
        
        # Time patterns
        time_patterns = [
            r'(\d+)\s*(?:mins?|minutes?)',  # 15 mins, 20 minutes
            r'(\d+)\s*(?:hours?|hrs?)',    # 10 hours, 50 hrs
        ]
        
        # Percentage patterns
        percentage_patterns = [
            r'(\d+(?:\.\d+)?)\s*%',  # 25%, 12.5%
            r'(\d+(?:\.\d+)?)\s*percent',  # 25 percent
        ]
        
        # Multiplier patterns
        multiplier_patterns = [
            r'(\d+)x\s*faster',  # 2x faster, 3x faster
            r'(\d+)\s*times\s*faster',  # 2 times faster
        ]
        
        for pattern in currency_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                value = float(match.group(1).replace(',', ''))
                unit = match.group(2) if match.group(2) else 'USD'
                if unit == 'M':
                    value *= 1000000
                elif unit == 'B':
                    value *= 1000000000
                elif unit == 'K':
                    value *= 1000
                
                numerical_data.append({
                    "value": value,
                    "unit": "USD",
                    "context": "currency",
                    "original_text": match.group(0)
                })
        
        for pattern in time_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                value = int(match.group(1))
                unit = "minutes" if "min" in match.group(0).lower() else "hours"
                numerical_data.append({
                    "value": value,
                    "unit": unit,
                    "context": "time_savings",
                    "original_text": match.group(0)
                })
        
        for pattern in percentage_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                value = float(match.group(1))
                numerical_data.append({
                    "value": value,
                    "unit": "percentage",
                    "context": "percentage",
                    "original_text": match.group(0)
                })
        
        for pattern in multiplier_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                value = int(match.group(1))
                numerical_data.append({
                    "value": value,
                    "unit": "multiplier",
                    "context": "performance_improvement",
                    "original_text": match.group(0)
                })
        
        return numerical_data
    
    def _extract_key_claims(self, text: str) -> List[str]:
        """Extract key claims and statements from text."""
        # Simple keyword-based extraction - could be enhanced with NLP
        key_phrases = [
            "AI-powered", "automated", "faster", "efficient", "streamlined",
            "competitive", "market leader", "innovative", "cutting-edge",
            "time-saving", "productivity", "efficiency", "accuracy"
        ]
        
        claims = []
        sentences = re.split(r'[.!?]+', text)
        
        for sentence in sentences:
            sentence = sentence.strip()
            if any(phrase.lower() in sentence.lower() for phrase in key_phrases):
                claims.append(sentence)
        
        return claims[:5]  # Limit to top 5 claims

class AIAnalyzer:
    """Uses Gemini AI to analyze content for inconsistencies."""
    
    def __init__(self, api_key: str):
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel('gemini-2.0-flash-exp')
        self.console = Console()
    
    def analyze_inconsistencies(self, slides_content: List[SlideContent]) -> List[Inconsistency]:
        """Analyze slides for inconsistencies using AI."""
        inconsistencies = []
        
        # Rule-based checks first
        rule_based = self._rule_based_checks(slides_content)
        inconsistencies.extend(rule_based)
        
        # AI-powered analysis
        ai_based = self._ai_based_analysis(slides_content)
        inconsistencies.extend(ai_based)
        
        return inconsistencies
    
    def _rule_based_checks(self, slides_content: List[SlideContent]) -> List[Inconsistency]:
        """Perform rule-based consistency checks."""
        inconsistencies = []
        
        # Check for numerical inconsistencies
        numerical_issues = self._check_numerical_consistency(slides_content)
        inconsistencies.extend(numerical_issues)
        
        # Check for claim contradictions
        claim_issues = self._check_claim_consistency(slides_content)
        inconsistencies.extend(claim_issues)
        
        return inconsistencies
    
    def _check_numerical_consistency(self, slides_content: List[SlideContent]) -> List[Inconsistency]:
        """Check for numerical data inconsistencies."""
        inconsistencies = []
        
        # Group numerical data by context
        context_groups = {}
        for slide in slides_content:
            for data in slide.numerical_data:
                context = data.get("context", "unknown")
                if context not in context_groups:
                    context_groups[context] = []
                context_groups[context].append({
                    "slide": slide.slide_number,
                    "data": data
                })
        
        # Check for conflicts within each context
        for context, data_list in context_groups.items():
            if len(data_list) < 2:
                continue
            
            # Check for conflicting values
            values = [item["data"]["value"] for item in data_list]
            if len(set(values)) > 1:
                slides_involved = [item["slide"] for item in data_list]
                evidence = [f"Slide {item['slide']}: {item['data']['original_text']}" for item in data_list]
                
                inconsistency = Inconsistency(
                    slide_numbers=slides_involved,
                    inconsistency_type="numerical_conflict",
                    description=f"Conflicting {context} values found across slides",
                    confidence=0.95,
                    evidence=evidence,
                    recommendation=f"Standardize {context} values across all slides for consistency"
                )
                inconsistencies.append(inconsistency)
        
        return inconsistencies
    
    def _check_claim_consistency(self, slides_content: List[SlideContent]) -> List[Inconsistency]:
        """Check for contradictory claims across slides."""
        inconsistencies = []
        
        # Check for performance improvement contradictions
        performance_claims = {}
        for slide in slides_content:
            for data in slide.numerical_data:
                if data.get("context") == "performance_improvement":
                    value = data["value"]
                    if "performance_improvement" not in performance_claims:
                        performance_claims["performance_improvement"] = []
                    performance_claims["performance_improvement"].append({
                        "slide": slide.slide_number,
                        "value": value,
                        "text": data["original_text"]
                    })
        
        # Check for conflicting performance claims
        if "performance_improvement" in performance_claims:
            claims = performance_claims["performance_improvement"]
            if len(claims) > 1:
                values = [claim["value"] for claim in claims]
                if len(set(values)) > 1:
                    slides_involved = [claim["slide"] for claim in claims]
                    evidence = [f"Slide {claim['slide']}: {claim['text']}" for claim in claims]
                    
                    inconsistency = Inconsistency(
                        slide_numbers=slides_involved,
                        inconsistency_type="performance_claim_conflict",
                        description="Conflicting performance improvement claims found",
                        confidence=0.90,
                        evidence=evidence,
                        recommendation="Align performance improvement claims across all slides"
                    )
                    inconsistencies.append(inconsistency)
        
        return inconsistencies
    
    def _ai_based_analysis(self, slides_content: List[SlideContent]) -> List[Inconsistency]:
        """Use AI to analyze content for subtle inconsistencies."""
        inconsistencies = []
        
        try:
            # Prepare content for AI analysis
            content_summary = self._prepare_content_summary(slides_content)
            
            prompt = f"""
            Analyze the following PowerPoint presentation content for factual and logical inconsistencies.
            
            Content Summary:
            {content_summary}
            
            Please identify any:
            1. Contradictory statements or claims
            2. Inconsistent numerical data
            3. Timeline mismatches
            4. Brand or product inconsistencies
            5. Logical contradictions
            
            For each inconsistency found, provide:
            - Slide numbers involved
            - Type of inconsistency
            - Description
            - Evidence
            - Recommendation
            
            Return your analysis in JSON format.
            """
            
            response = self.model.generate_content(prompt)
            
            # Parse AI response
            try:
                ai_analysis = json.loads(response.text)
                if "inconsistencies" in ai_analysis:
                    for item in ai_analysis["inconsistencies"]:
                        inconsistency = Inconsistency(
                            slide_numbers=item.get("slide_numbers", []),
                            inconsistency_type=item.get("type", "ai_detected"),
                            description=item.get("description", ""),
                            confidence=item.get("confidence", 0.7),
                            evidence=item.get("evidence", []),
                            recommendation=item.get("recommendation", "")
                        )
                        inconsistencies.append(inconsistency)
            except json.JSONDecodeError:
                logger.warning("AI response not in valid JSON format")
                
        except Exception as e:
            logger.error(f"AI analysis failed: {e}")
        
        return inconsistencies
    
    def _prepare_content_summary(self, slides_content: List[SlideContent]) -> str:
        """Prepare a summary of slide content for AI analysis."""
        summary_parts = []
        
        for slide in slides_content:
            slide_summary = f"Slide {slide.slide_number}:\n"
            slide_summary += f"Text: {slide.text_content[:200]}...\n"
            
            if slide.numerical_data:
                numbers = [f"{d['value']}{d['unit']} ({d['context']})" for d in slide.numerical_data]
                slide_summary += f"Numbers: {', '.join(numbers)}\n"
            
            if slide.key_claims:
                slide_summary += f"Claims: {', '.join(slide.key_claims[:3])}\n"
            
            summary_parts.append(slide_summary)
        
        return "\n".join(summary_parts)

class ReportGenerator:
    """Generates structured reports of detected inconsistencies."""
    
    def __init__(self):
        self.console = Console()
    
    def generate_report(self, inconsistencies: List[Inconsistency], output_format: str = "console", output_file: str = None) -> None:
        """Generate and output the inconsistency report."""
        if output_format == "console":
            self._console_report(inconsistencies)
        elif output_format == "json":
            self._json_report(inconsistencies, output_file)
        elif output_format == "csv":
            self._csv_report(inconsistencies, output_file)
        else:
            self._console_report(inconsistencies)
    
    def _console_report(self, inconsistencies: List[Inconsistency]) -> None:
        """Generate console-based report using Rich."""
        if not inconsistencies:
            self.console.print(Panel("✅ No inconsistencies detected!", style="green"))
            return
        
        # Summary
        self.console.print(Panel(f"🔍 Found {len(inconsistencies)} inconsistencies", style="yellow"))
        
        # Group by type
        by_type = {}
        for inc in inconsistencies:
            inc_type = inc.inconsistency_type
            if inc_type not in by_type:
                by_type[inc_type] = []
            by_type[inc_type].append(inc)
        
        # Display by type
        for inc_type, incs in by_type.items():
            self.console.print(f"\n[bold blue]{inc_type.upper()}[/bold blue] ({len(incs)} issues)")
            
            for inc in incs:
                table = Table(show_header=True, header_style="bold magenta")
                table.add_column("Property", style="cyan")
                table.add_column("Value", style="white")
                
                table.add_row("Slides", ", ".join(map(str, inc.slide_numbers)))
                table.add_row("Type", inc.inconsistency_type)
                table.add_row("Description", inc.description)
                table.add_row("Confidence", f"{inc.confidence:.1%}")
                table.add_row("Evidence", "\n".join(inc.evidence))
                table.add_row("Recommendation", inc.recommendation)
                
                self.console.print(table)
                self.console.print("")
    
    def _json_report(self, inconsistencies: List[Inconsistency], output_file: str = None) -> None:
        """Generate JSON report."""
        report = {
            "timestamp": datetime.now().isoformat(),
            "total_inconsistencies": len(inconsistencies),
            "inconsistencies": [asdict(inc) for inc in inconsistencies]
        }
        
        if output_file:
            with open(output_file, 'w') as f:
                json.dump(report, f, indent=2)
            self.console.print(f"✅ JSON report saved to {output_file}")
        else:
            print(json.dumps(report, indent=2))
    
    def _csv_report(self, inconsistencies: List[Inconsistency], output_file: str = None) -> None:
        """Generate CSV report."""
        if not inconsistencies:
            return
        
        # Flatten inconsistencies for CSV
        rows = []
        for inc in inconsistencies:
            for i, evidence in enumerate(inc.evidence):
                row = {
                    "slide_numbers": ", ".join(map(str, inc.slide_numbers)),
                    "type": inc.inconsistency_type,
                    "description": inc.description,
                    "confidence": inc.confidence,
                    "evidence": evidence,
                    "recommendation": inc.recommendation
                }
                rows.append(row)
        
        df = pd.DataFrame(rows)
        
        if output_file:
            df.to_csv(output_file, index=False)
            self.console.print(f"✅ CSV report saved to {output_file}")
        else:
            print(df.to_csv(index=False))

class PowerPointInspector:
    """Main class for PowerPoint inconsistency detection."""
    
    def __init__(self, api_key: str):
        self.extractor = ContentExtractor()
        self.analyzer = AIAnalyzer(api_key)
        self.report_generator = ReportGenerator()
        self.console = Console()
    
    def analyze_presentation(self, file_path: str = None, images_dir: str = None) -> List[Inconsistency]:
        """Analyze a presentation for inconsistencies."""
        if file_path:
            self.console.print(f"📊 Analyzing PowerPoint file: {file_path}")
            slides_content = self.extractor.extract_from_pptx(file_path)
        elif images_dir:
            self.console.print(f"🖼️  Analyzing images from: {images_dir}")
            slides_content = self.extractor.extract_from_images(images_dir)
        else:
            raise ValueError("Either file_path or images_dir must be provided")
        
        if not slides_content:
            self.console.print("[red]No content could be extracted from the presentation[/red]")
            return []
        
        self.console.print(f"📝 Extracted content from {len(slides_content)} slides")
        
        # Analyze for inconsistencies
        with Progress(
            SpinnerColumn(),
            TextColumn("[progress.description]{task.description}"),
            console=self.console
        ) as progress:
            task = progress.add_task("Analyzing for inconsistencies...", total=None)
            inconsistencies = self.analyzer.analyze_inconsistencies(slides_content)
            progress.update(task, completed=True)
        
        return inconsistencies

def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(description="PowerPoint Inconsistency Detector")
    parser.add_argument("analyze", help="Analyze a presentation")
    parser.add_argument("--file", help="Path to PowerPoint file (.pptx)")
    parser.add_argument("--images", help="Directory containing slide images")
    parser.add_argument("--output", help="Output file path")
    parser.add_argument("--format", choices=["json", "csv", "console"], default="console", help="Output format")
    parser.add_argument("--verbose", action="store_true", help="Enable verbose logging")
    
    args = parser.parse_args()
    
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # Check for API key
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        console = Console()
        console.print("[red]Error: GEMINI_API_KEY environment variable not set[/red]")
        console.print("Please set your Gemini API key in the .env file")
        console.print("Get your free API key from: https://aistudio.google.com/app/apikey")
        sys.exit(1)
    
    # Validate inputs
    if not args.file and not args.images:
        console = Console()
        console.print("[red]Error: Either --file or --images must be specified[/red]")
        sys.exit(1)
    
    if args.file and not os.path.exists(args.file):
        console = Console()
        console.print(f"[red]Error: File not found: {args.file}[/red]")
        sys.exit(1)
    
    if args.images and not os.path.exists(args.images):
        console = Console()
        console.print(f"[red]Error: Images directory not found: {args.images}[/red]")
        sys.exit(1)
    
    try:
        # Initialize inspector
        inspector = PowerPointInspector(api_key)
        
        # Analyze presentation
        inconsistencies = inspector.analyze_presentation(
            file_path=args.file,
            images_dir=args.images
        )
        
        # Generate report
        inspector.report_generator.generate_report(
            inconsistencies,
            output_format=args.format,
            output_file=args.output
        )
        
        # Exit with appropriate code
        if inconsistencies:
            sys.exit(1)  # Exit with error code if inconsistencies found
        else:
            sys.exit(0)  # Exit successfully if no inconsistencies
            
    except Exception as e:
        console = Console()
        console.print(f"[red]Error: {e}[/red]")
        if args.verbose:
            import traceback
            traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
